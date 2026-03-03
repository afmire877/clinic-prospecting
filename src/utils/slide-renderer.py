#!/usr/bin/env python3
"""
Fallback slide renderer using python-pptx + Pillow.
Renders PPTX slides to PNG images when LibreOffice is unavailable.

Usage: python3 slide-renderer.py <input.pptx> <output-dir> [width] [height]
"""
import sys
import os

def render_slides(input_path, output_dir, width=1920, height=1080):
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.enum.text import PP_ALIGN
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        print(f"ERROR: Missing dependency: {e}", file=sys.stderr)
        print("Install with: pip install python-pptx Pillow", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(input_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Calculate scaling
    scale_x = width / slide_w
    scale_y = height / slide_h
    scale = min(scale_x, scale_y)

    os.makedirs(output_dir, exist_ok=True)

    # Try to find a usable font
    font_path = None
    font_candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "C:\\Windows\\Fonts\\arial.ttf",
    ]
    for fp in font_candidates:
        if os.path.exists(fp):
            font_path = fp
            break

    for idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (width, height), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Render shapes
        for shape in slide.shapes:
            x = int(shape.left * scale)
            y = int(shape.top * scale)
            w = int(shape.width * scale)
            h = int(shape.height * scale)

            # Draw filled rectangle for shapes with fill
            if hasattr(shape, "fill") and shape.fill and shape.fill.type is not None:
                try:
                    rgb = shape.fill.fore_color.rgb
                    color = (rgb[0], rgb[1], rgb[2])
                    draw.rectangle([x, y, x + w, y + h], fill=color)
                except Exception:
                    pass

            # Draw text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    # Get font size
                    font_size = 24  # default
                    font_color = (0, 0, 0)  # default black
                    for run in paragraph.runs:
                        if run.font.size:
                            font_size = int(run.font.size * scale / 12700)
                        try:
                            if run.font.color and run.font.color.type is not None:
                                rgb = run.font.color.rgb
                                if rgb:
                                    font_color = (rgb[0], rgb[1], rgb[2])
                        except (AttributeError, TypeError):
                            pass

                    font_size = max(12, min(font_size, 200))

                    try:
                        if font_path:
                            font = ImageFont.truetype(font_path, font_size)
                        else:
                            font = ImageFont.load_default()
                    except Exception:
                        font = ImageFont.load_default()

                    # Get text bounding box for centering
                    bbox = draw.textbbox((0, 0), text, font=font)
                    tw = bbox[2] - bbox[0]
                    th = bbox[3] - bbox[1]

                    # Center text in the text box
                    tx = x + (w - tw) // 2
                    ty = y + (h - th) // 2

                    draw.text((tx, ty), text, fill=font_color, font=font)

        out_path = os.path.join(output_dir, f"slide_{idx + 1:03d}.png")
        img.save(out_path, "PNG")

    print(f"Rendered {len(prs.slides)} slides to {output_dir}")
    return len(prs.slides)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <output-dir> [width] [height]")
        sys.exit(1)

    input_path = sys.argv[1]
    output_dir = sys.argv[2]
    width = int(sys.argv[3]) if len(sys.argv) > 3 else 1920
    height = int(sys.argv[4]) if len(sys.argv) > 4 else 1080

    if not os.path.exists(input_path):
        print(f"Error: File not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    render_slides(input_path, output_dir, width, height)
